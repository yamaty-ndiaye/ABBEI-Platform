import os
import shutil
from pathlib import Path
import pdfplumber
import fitz
from docx import Document as DocxDocument
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_ollama import OllamaEmbeddings
from langchain_core.documents import Document
from sentence_transformers import CrossEncoder
from rank_bm25 import BM25Okapi
import anthropic
from app.agents.dropbox_connector import (
    get_client, lister_fichiers_dropbox, telecharger_fichier, generer_lien
)
from app.core.config import (
    OLLAMA_HOST, MISTRAL_MODEL, EMBEDDING_MODEL,
    DROPBOX_CHANTIERS_PATH, CHROMA_PATH,
    CHUNK_SIZE, CHUNK_OVERLAP, SYSTEM_PROMPT,
    ANTHROPIC_API_KEY, CLAUDE_MODEL
)


# ============================================
# Extraction de texte
# ============================================

def extraire_texte(tmp_path: str, extension: str) -> str:
    try:
        if extension == '.pdf':
            texte = ""

            # Essai 1 : pdfplumber
            try:
                with pdfplumber.open(tmp_path) as pdf:
                    for page in pdf.pages:
                        t = page.extract_text()
                        if t:
                            texte += t + "\n"
            except:
                pass

            # Essai 2 : OCR sur toutes les pages
            try:
                from pdf2image import convert_from_path
                import pytesseract
                pages_images = convert_from_path(tmp_path, dpi=150)
                for page_img in pages_images:
                    ocr_texte = pytesseract.image_to_string(page_img, lang='fra')
                    if ocr_texte.strip():
                        texte += ocr_texte + "\n"
            except:
                pass

            # Fallback PyMuPDF
            if not texte.strip():
                try:
                    doc = fitz.open(tmp_path)
                    for page in doc:
                        texte += page.get_text()
                except:
                    pass

            return texte

        elif extension in ['.docx', '.doc']:
            try:
                doc = DocxDocument(tmp_path)
                texte = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
                if texte.strip():
                    return texte
                import docx2txt
                return docx2txt.process(tmp_path) or ""
            except:
                try:
                    import docx2txt
                    return docx2txt.process(tmp_path) or ""
                except:
                    return ""

        elif extension in ['.xlsx', '.xls']:
            try:
                import openpyxl
                wb = openpyxl.load_workbook(tmp_path)
                texte = ""
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        ligne = " ".join([str(c) for c in row if c is not None])
                        if ligne.strip():
                            texte += ligne + "\n"
                return texte
            except:
                return ""

        elif extension == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(tmp_path)
                texte = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texte += shape.text + "\n"
                return texte
            except:
                return ""

        else:
            return ""
    except Exception:
        return ""


# ============================================
# Détection du type de document
# ============================================

def detecter_type(nom: str) -> str:
    nom_lower = nom.lower()
    if any(x in nom_lower for x in ['facture', 'f150206', 'f240', 'f260']):
        return "facture"
    elif any(x in nom_lower for x in ['bc.', 'bt.', 'otms', 'otsol', 'otpeint', 'otdgs']):
        return "bon_commande"
    elif any(x in nom_lower for x in ['o02', 'o03', 'o04', 'o05', 'ar-09', 'ar-10', 'ar-13']):
        return "administratif"
    else:
        return "autre"


def detecter_filtre(question: str) -> dict:
    q = question.lower()
    if any(x in q for x in ['facture', 'montant', 'ttc', 'ht', 'tva', 'paiement']):
        return {"type": "facture"}
    elif any(x in q for x in ['bon de commande', 'bc', 'intervention', 'bt']):
        return {"type": "bon_commande"}
    return {}


# ============================================
# Chunking intelligent par type de document
# ============================================

def chunker_hybrid(doc: Document) -> list:
    if len(doc.page_content) < 2000:
        return [doc]
    splitter = RecursiveCharacterTextSplitter(
        separators=["\n\n", "\n", " "],
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP
    )
    return splitter.split_documents([doc])


def chunker_paragraphe(doc: Document) -> list:
    splitter = RecursiveCharacterTextSplitter(
        separators=["\n\n", "\n"],
        chunk_size=1500,
        chunk_overlap=150
    )
    return splitter.split_documents([doc])


def chunker_simple(doc: Document) -> list:
    if len(doc.page_content) < 2000:
        return [doc]
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=50
    )
    return splitter.split_documents([doc])


def chunker_par_type(doc: Document) -> list:
    type_doc = doc.metadata.get("type", "autre")
    if type_doc == "facture":
        return chunker_hybrid(doc)
    elif type_doc == "administratif":
        return chunker_paragraphe(doc)
    elif type_doc == "bon_commande":
        return chunker_simple(doc)
    else:
        return chunker_simple(doc)


# ============================================
# Recherche hybride BM25 + Vector
# ============================================

def recherche_bm25(chunks_tous: list, question: str, k: int = 20) -> list:
    corpus = [chunk.page_content.lower().split() for chunk in chunks_tous]
    bm25 = BM25Okapi(corpus)
    tokens_question = question.lower().split()
    scores = bm25.get_scores(tokens_question)
    indices_tries = sorted(range(len(scores)), key=lambda i: scores[i], reverse=True)[:k]
    return [chunks_tous[i] for i in indices_tries]


def recherche_hybride(vectorstore, question: str, k: int = 20) -> list:
    chunks_vector = vectorstore.similarity_search(question, k=k)
    chunks_bm25 = recherche_bm25(chunks_vector, question, k=k)

    vus = set()
    chunks_fusionnes = []
    for chunk_v, chunk_b in zip(chunks_vector, chunks_bm25):
        for chunk in [chunk_v, chunk_b]:
            nom = chunk.metadata.get("nom", "")
            if nom not in vus:
                chunks_fusionnes.append(chunk)
                vus.add(nom)

    return chunks_fusionnes


# ============================================
# Indexation
# ============================================

def indexer_documents(fichiers_liste: list = None, chemin_dropbox: str = None, reset: bool = False):
    chemin = chemin_dropbox or DROPBOX_CHANTIERS_PATH

    if reset and os.path.exists(CHROMA_PATH):
        shutil.rmtree(CHROMA_PATH)
        print("🗑️ ChromaDB réinitialisée")

    dbx = get_client()

    # Utiliser la liste fournie ou lister depuis Dropbox
    if fichiers_liste:
        fichiers = fichiers_liste
    else:
        print(f"🔍 Indexation depuis Dropbox : {chemin}")
        fichiers = lister_fichiers_dropbox(dbx, chemin)

    print(f"✅ {len(fichiers)} fichiers à indexer")

    print("📄 Extraction du texte...")
    documents = []
    erreurs = 0

    for i, fichier in enumerate(fichiers):
        tmp_path = None
        try:
            tmp_path = telecharger_fichier(dbx, fichier["chemin"])
            texte = extraire_texte(tmp_path, fichier["extension"])

            if texte.strip():
                documents.append(Document(
                    page_content=texte,
                    metadata={
                        "source": fichier["chemin"],
                        "nom": fichier["nom"],
                        "extension": fichier["extension"],
                        "type": detecter_type(fichier["nom"])
                    }
                ))
            else:
                erreurs += 1

        except Exception as e:
            erreurs += 1
            print(f"  ❌ {fichier['nom']}: {e}")
        finally:
            if tmp_path and os.path.exists(tmp_path):
                os.unlink(tmp_path)

        if i % 5 == 0:
            print(f"  → {i}/{len(fichiers)} traités ({len(documents)} lisibles)...")

    print(f"✅ {len(documents)} documents extraits ({erreurs} erreurs)")

    if not documents:
        print("❌ Aucun document lisible")
        return None

    print("✂️  Découpage en chunks...")
    chunks = []
    for doc in documents:
        chunks.extend(chunker_par_type(doc))

    print(f"✅ {len(chunks)} chunks créés")

    print("🧠 Création des embeddings...")
    embeddings = OllamaEmbeddings(model=EMBEDDING_MODEL, base_url=OLLAMA_HOST)
    vectorstore = Chroma.from_documents(
        documents=chunks,
        embedding=embeddings,
        persist_directory=CHROMA_PATH
    )
    print(f"✅ Indexation terminée — {len(chunks)} chunks dans ChromaDB")
    return vectorstore


# ============================================
# Chargement du vectorstore
# ============================================

def charger_vectorstore():
    embeddings = OllamaEmbeddings(model=EMBEDDING_MODEL, base_url=OLLAMA_HOST)
    return Chroma(persist_directory=CHROMA_PATH, embedding_function=embeddings)


# ============================================
# Reranker (singleton)
# ============================================

_reranker = None

def get_reranker():
    global _reranker
    if _reranker is None:
        _reranker = CrossEncoder("cross-encoder/ms-marco-MiniLM-L-6-v2")
    return _reranker


# ============================================
# Réponse RAG
# ============================================

def repondre(question: str) -> dict:
    vectorstore = charger_vectorstore()

    filtre = detecter_filtre(question)
    if filtre:
        chunks = vectorstore.similarity_search(question, k=50, filter=filtre)
        if not chunks:
            chunks = recherche_hybride(vectorstore, question, k=50)
    else:
        chunks = recherche_hybride(vectorstore, question, k=50)

    if not chunks:
        return {
            "question": question,
            "reponse": "Aucun document indexé pour répondre à cette question.",
            "sources": []
        }

    chunks_dedupliques = []
    sources_vues = set()
    for chunk in chunks:
        nom = chunk.metadata.get("nom", "")
        if nom not in sources_vues:
            chunks_dedupliques.append(chunk)
            sources_vues.add(nom)

    reranker = get_reranker()
    scores = reranker.predict(
        [(question, chunk.page_content) for chunk in chunks_dedupliques]
    )

    chunks_tries = [chunk for _, chunk in
                    sorted(zip(scores, chunks_dedupliques), key=lambda x: x[0], reverse=True)][:3]

    contexte = "\n\n".join([c.page_content for c in chunks_tries])

    prompt_text = f"""{SYSTEM_PROMPT}

Contexte extrait des documents ABBEI :
{contexte}

Question : {question}

Réponse :"""

    client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
    message = client.messages.create(
        model=CLAUDE_MODEL,
        max_tokens=1024,
        messages=[{"role": "user", "content": prompt_text}]
    )
    reponse = message.content[0].text

    dbx = get_client()
    sources = []
    vus = set()
    for chunk in chunks_tries:
        chemin = chunk.metadata.get("source", "")
        if chemin in vus:
            continue
        vus.add(chemin)
        nom = chunk.metadata.get("nom", chemin)
        lien = generer_lien(dbx, chemin) if chemin.lower().startswith("/abbei") else None
        sources.append({
            "nom": nom,
            "chemin": chemin,
            "lien": lien
        })

    return {
        "question": question,
        "reponse": reponse,
        "sources": sources
    }


# ============================================
# Test sur 20 fichiers
# ============================================

if __name__ == "__main__":
    dbx = get_client()
    tous_fichiers = lister_fichiers_dropbox(
        dbx,
        "/abbei/ChantiersABBEI/HABITAT 76/Y-PE010Y-H76-MarchéEntretien"
    )

    # Sélectionner 30 fichiers représentatifs
    factures = [f for f in tous_fichiers if any(x in f['nom'].lower() for x in ['f150206', 'f240', 'f260', 'facture'])]
    bons_commande = [f for f in tous_fichiers if any(x in f['nom'].lower() for x in [' bt', ' bc', 'otms', 'otsol', 'otpeint', 'otdgs'])]
    administratifs = [f for f in tous_fichiers if any(x in f['nom'].lower() for x in ['o02', 'o03', 'ccap', 'cctp', 'marche'])]
    autres = [f for f in tous_fichiers if f not in factures + bons_commande + administratifs]

    # Prendre un mix représentatif
    fichiers_test = (
        factures[:5] +        # 5 factures
        bons_commande[:10] +  # 10 bons de commande
        administratifs[:5] +  # 5 administratifs
        autres[:10]           # 10 autres
    )[:30]

    print(f"📋 {len(fichiers_test)} fichiers sélectionnés :")
    print(f"  Factures      : {len([f for f in fichiers_test if detecter_type(f['nom']) == 'facture'])}")
    print(f"  Bons commande : {len([f for f in fichiers_test if detecter_type(f['nom']) == 'bon_commande'])}")
    print(f"  Administratifs: {len([f for f in fichiers_test if detecter_type(f['nom']) == 'administratif'])}")
    print(f"  Autres        : {len([f for f in fichiers_test if detecter_type(f['nom']) == 'autre'])}")

    if os.path.exists(CHROMA_PATH):
        shutil.rmtree(CHROMA_PATH)

    indexer_documents(fichiers_liste=fichiers_test, reset=False)